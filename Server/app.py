from flask import Flask, request, jsonify
from flask_cors import CORS
import pytesseract
from PIL import Image
import logging
import nltk
from nltk.tokenize import sent_tokenize
import random
import jwt
import datetime
import os
from dotenv import load_dotenv
import bcrypt
import re
import PyPDF2
import google.generativeai as genai
from pptx import Presentation
import json

# Explicitly set Tesseract path
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()

app = Flask(__name__)
CORS(app, resources={r"/*": {
    "origins": "http://localhost:3000",
    "methods": ["GET", "POST", "OPTIONS"],
    "allow_headers": ["Content-Type", "Authorization"],
    "supports_credentials": True
}})

# Initialize Gemini API
try:
    genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
    gemini_model = genai.GenerativeModel('gemini-1.5-flash')
    logger.info("Gemini API initialized successfully")
except Exception as e:
    logger.error(f"Failed to initialize Gemini API: {str(e)}")
    raise Exception(f"Gemini API initialization failed: {str(e)}")

# Download NLTK resources
try:
    nltk.download('punkt', quiet=True)
    nltk.download('punkt_tab', quiet=True)
    logger.info("NLTK resources (punkt, punkt_tab) downloaded successfully")
except Exception as e:
    logger.error(f"Failed to download NLTK resources: {str(e)}")
    raise Exception(f"NLTK resource download failed: {str(e)}")

# JWT Secret
SECRET_KEY = os.getenv('SECRET_KEY', 'your-secret-key')
if SECRET_KEY == 'your-secret-key':
    logger.warning("Using default SECRET_KEY. Set a secure key in .env for production.")

# Mock databases
users = {}  # {username: hashed_password}
quiz_sessions = {}  # {user_id: [questions]}

# Input sanitization
def sanitize_input(text):
    if not text:
        return ""
    return re.sub(r'[^\w\s.,!?\'"]', '', text)[:1000]  # Increased limit, allow basic punctuation

# Authentication middleware
def authenticate_token(token):
    if not token:
        logger.warning("No token provided in request")
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
        logger.debug(f"Token validated for user: {payload['user_id']}")
        return payload['user_id']
    except jwt.ExpiredSignatureError:
        logger.warning("Token expired")
        return None
    except jwt.InvalidTokenError:
        logger.warning("Invalid token")
        return None

@app.route('/register', methods=['POST'])
def register():
    try:
        data = request.get_json(silent=True)
        if not data:
            logger.error("No JSON data provided in register request")
            return jsonify({'error': 'Invalid or missing JSON data'}), 400

        username = sanitize_input(data.get('username'))
        password = data.get('password')

        if not username or not password:
            logger.error("Missing username or password in register request")
            return jsonify({'error': 'Username and password are required'}), 400

        if len(username) < 3 or len(password) < 6:
            logger.error(f"Invalid input: username={username}, password length={len(password)}")
            return jsonify({'error': 'Username must be at least 3 characters and password at least 6 characters'}), 400

        if username in users:
            logger.warning(f"Attempt to register existing user: {username}")
            return jsonify({'error': 'User already exists'}), 400

        hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
        users[username] = hashed_password.decode('utf-8')
        logger.info(f"User registered successfully: {username}")

        return jsonify({'message': 'User registered successfully'}), 201
    except Exception as e:
        logger.error(f"Registration failed: {str(e)}")
        return jsonify({'error': f'Registration failed: {str(e)}'}), 500

@app.route('/login', methods=['POST'])
def login():
    try:
        data = request.get_json(silent=True)
        if not data:
            logger.error("No JSON data provided in login request")
            return jsonify({'error': 'Invalid or missing JSON data'}), 400

        username = sanitize_input(data.get('username'))
        password = data.get('password')

        if not username or not password:
            logger.error("Missing username or password in login request")
            return jsonify({'error': 'Username and password are required'}), 400

        stored_password = users.get(username)
        if not stored_password or not bcrypt.checkpw(password.encode('utf-8'), stored_password.encode('utf-8')):
            logger.warning(f"Invalid login attempt for user: {username}")
            return jsonify({'error': 'Invalid credentials'}), 401

        token = jwt.encode({
            'user_id': username,
            'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=24)
        }, SECRET_KEY, algorithm="HS256")
        logger.info(f"User logged in successfully: {username}")

        return jsonify({'token': token}), 200
    except Exception as e:
        logger.error(f"Login failed: {str(e)}")
        return jsonify({'error': f'Login failed: {str(e)}'}), 500

@app.route('/ocr', methods=['POST'])
def ocr():
    logger.debug("Received request to /ocr endpoint")
    auth_header = request.headers.get('Authorization')
    logger.debug(f"Authorization header: {auth_header}")
    token = auth_header.split(" ")[1] if auth_header and len(auth_header.split(" ")) > 1 else None
    if not authenticate_token(token):
        logger.error("Unauthorized access to OCR endpoint")
        return jsonify({'error': 'Unauthorized'}), 401
    try:
        logger.debug("Checking for file in request")
        if 'file' not in request.files:
            logger.error("No file provided in OCR request")
            return jsonify({'error': 'No file provided'}), 400
        file = request.files['file']
        filename = file.filename.lower()
        logger.debug(f"Processing file: {filename}")

        text = ""
        if filename.endswith('.pdf'):
            logger.debug("Processing PDF file")
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif filename.endswith(('.jpg', '.jpeg', '.png')):
            logger.debug("Processing image file")
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
        elif filename.endswith(('.ppt', '.pptx')):
            logger.debug("Processing PowerPoint file")
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            logger.error("Unsupported file format")
            return jsonify({'error': 'Unsupported file format. Use JPG, PNG, PDF, PPT, or PPTX'}), 400

        logger.debug(f"Extracted text length: {len(text)}")
        if not text.strip():
            logger.warning("No text extracted from file")
            return jsonify({'error': 'No text extracted from file'}), 400

        # Summarize using Gemini API
        logger.debug("Generating summary with Gemini API")
        prompt = f"Summarize the following text in a clear and understandable way (max 150 words):\n\n{text}"
        try:
            response = gemini_model.generate_content(prompt)
            summary = response.text.strip()
            if not summary:
                logger.warning("Gemini API returned empty summary")
                summary = "No summary generated."
        except Exception as e:
            logger.error(f"Gemini summarization failed: {str(e)}")
            summary = "Failed to generate summary."
        logger.info("OCR and summarization completed successfully")
        return jsonify({'text': text, 'summary': summary})
    except Exception as e:
        logger.error(f"OCR processing failed: {str(e)}")
        return jsonify({'error': f'OCR processing failed: {str(e)}'}), 500

@app.route('/ask', methods=['POST'])
def ask():
    auth_header = request.headers.get('Authorization')
    token = auth_header.split(" ")[1] if auth_header and len(auth_header.split(" ")) > 1 else None
    if not authenticate_token(token):
        logger.error("Unauthorized access to ask endpoint")
        return jsonify({'error': 'Unauthorized'}), 401
    try:
        data = request.get_json(silent=True)
        if not data:
            logger.error("No JSON data provided in ask request")
            return jsonify({'error': 'Invalid or missing JSON data'}), 400
        question = sanitize_input(data.get('question'))
        if not question:
            logger.error("Missing question in ask request")
            return jsonify({'error': 'Question is required'}), 400

        # Answer using Gemini API
        logger.debug("Answering question with Gemini API")
        prompt = f"Provide a clear and accurate answer to the following study-related question (max 200 words):\n\n{question}"
        try:
            response = gemini_model.generate_content(prompt)
            answer = response.text.strip()
            if not answer:
                logger.warning("Gemini API returned empty answer")
                answer = "No answer generated."
        except Exception as e:
            logger.error(f"Gemini answer generation failed: {str(e)}")
            answer = "Failed to generate answer."
        logger.info(f"Question answered: {question[:50]}...")
        return jsonify({'answer': answer})  # Fixed typo: 'error' to 'answer'
    except Exception as e:
        logger.error(f"Question processing failed: {str(e)}")
        return jsonify({'error': f'Question processing failed: {str(e)}'}), 500

@app.route('/quiz', methods=['POST'])
def quiz():
    auth_header = request.headers.get('Authorization')
    token = auth_header.split(" ")[1] if auth_header and len(auth_header.split(" ")) > 1 else None
    user_id = authenticate_token(token)
    if not user_id:
        logger.error("Unauthorized access to quiz endpoint")
        return jsonify({'error': 'Unauthorized'}), 401
    try:
        data = request.get_json(silent=True)
        if not data:
            logger.error("No JSON data provided in quiz request")
            return jsonify({'error': 'Invalid or missing JSON data'}), 400
        context = sanitize_input(data.get('context'))
        user_answers = data.get('answers', [])  # List of {question_id, selected_option}
        if not context:
            logger.error("Missing context in quiz request")
            return jsonify({'error': 'Context or topic is required'}), 400

        # Generate questions if no answers provided (new quiz session)
        if not user_answers:
            logger.debug(f"Generating quiz questions for context: {context[:50]}...")
            sentences = sent_tokenize(context)
            if not sentences:
                logger.error("No sentences found in context after tokenization")
                return jsonify({'error': 'Invalid or too short context for quiz generation'}), 400

            bloom_levels = [
                {'level': 'Remember', 'prompt': 'What is'},
                {'level': 'Understand', 'prompt': 'Explain'},
                {'level': 'Apply', 'prompt': 'How would you apply'},
                {'level': 'Analyze', 'prompt': 'Compare and contrast'},
                {'level': 'Evaluate', 'prompt': 'Assess the importance of'}
            ]

            questions = []
            for i in range(5):
                level = random.choice(bloom_levels)
                key_sentence = random.choice(sentences)
                prompt = f"""
                Generate a {level['level']}-level multiple-choice question based on this context: {context}.
                Return a valid JSON object with:
                - question: string (1-2 sentences)
                - correct_answer: string (1 sentence)
                - distractors: array of 3 strings (1 sentence each)
                - explanation: string (1-2 sentences)
                Example:
                {{
                  "question": "What is photosynthesis?",
                  "correct_answer": "It is the process by which plants make food using sunlight.",
                  "distractors": ["It is the process of plant respiration.", "It is the process of seed germination.", "It is the process of root growth."],
                  "explanation": "Photosynthesis uses sunlight, CO2, and water to produce glucose and oxygen."
                }}
                Wrap the JSON in triple backticks (```json\n...\n```) to ensure proper formatting.
                """
                try:
                    response = gemini_model.generate_content(prompt)
                    response_text = response.text.strip()
                    logger.debug(f"Gemini raw response: {response_text[:500]}...")
                    # Extract JSON from code block if present
                    json_match = re.search(r'```json\n(.*?)\n```', response_text, re.DOTALL)
                    if json_match:
                        json_text = json_match.group(1)
                    else:
                        json_text = response_text
                    try:
                        quiz_data = json.loads(json_text)
                    except json.JSONDecodeError as e:
                        logger.error(f"Failed to parse JSON: {json_text[:200]}... Error: {str(e)}")
                        # Fallback question
                        quiz_data = {
                            'question': f"What is the main idea of the topic '{context}'?",
                            'correct_answer': f"It is a key concept related to {context}.",
                            'distractors': [
                                f"It is unrelated to {context}.",
                                f"It is a minor detail of {context}.",
                                f"It is the opposite of {context}."
                            ],
                            'explanation': f"This is a fallback question due to API parsing issues."
                        }
                    if not all(key in quiz_data for key in ['question', 'correct_answer', 'distractors', 'explanation']):
                        logger.error("Incomplete Gemini response structure")
                        return jsonify({'error': 'Invalid response structure from Gemini API'}), 500
                    if len(quiz_data['distractors']) != 3:
                        logger.error("Incorrect number of distractors")
                        return jsonify({'error': 'Invalid distractors in Gemini response'}), 500
                    options = [quiz_data['correct_answer']] + quiz_data['distractors']
                    random.shuffle(options)
                    questions.append({
                        'id': i + 1,
                        'question': quiz_data['question'],
                        'level': level['level'],
                        'options': options,
                        'correct_answer': quiz_data['correct_answer'],
                        'explanation': quiz_data['explanation']
                    })
                except Exception as e:
                    logger.error(f"Gemini API call failed: {str(e)}")
                    return jsonify({'error': f'Gemini API error: {str(e)}'}), 500
            # Store questions in mock database
            quiz_sessions[user_id] = questions
            logger.info(f"New quiz session generated with {len(questions)} questions for user: {user_id}")
            return jsonify({'questions': questions, 'time_per_question': 30})

        # Score submitted answers
        logger.debug(f"Scoring answers for user: {user_id}")
        if user_id not in quiz_sessions:
            logger.error(f"No quiz session found for user: {user_id}")
            return jsonify({'error': 'No active quiz session found'}), 400
        questions = quiz_sessions[user_id]
        logger.debug(f"Retrieved {len(questions)} questions for scoring")
        score = 0
        results = []
        for answer in user_answers:
            question_id = answer.get('question_id')
            selected_option = answer.get('selected_option')
            question_data = next((q for q in questions if q['id'] == question_id), None)
            if not question_data:
                logger.warning(f"Invalid question_id: {question_id}")
                continue
            is_correct = selected_option == question_data['correct_answer']
            if is_correct:
                score += 1
            results.append({
                'question_id': question_id,
                'selected_option': selected_option,
                'correct_answer': question_data['correct_answer'],
                'is_correct': is_correct,
                'explanation': question_data['explanation']
            })
        logger.info(f"Quiz scored: {score}/{len(user_answers)} for user: {user_id}")
        # Clear session after scoring
        quiz_sessions.pop(user_id, None)
        return jsonify({'score': score, 'total': len(user_answers), 'results': results})

    except Exception as e:
        logger.error(f"Quiz processing failed: {str(e)}")
        return jsonify({'error': f'Quiz processing failed: {str(e)}'}), 500

if __name__ == '__main__':
    try:
        logger.info("Starting Flask server...")
        app.run(debug=True, host='0.0.0.0', port=5000)
    except Exception as e:
        logger.error(f"Failed to start server: {str(e)}")
        raise